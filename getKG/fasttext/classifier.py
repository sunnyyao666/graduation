import fasttext
import jieba
import string
import os
from pptx import Presentation


def is_all_eng(strs):
    for i in strs:
        if i not in string.ascii_lowercase + string.ascii_uppercase:
            return False
    return True


classifier = fasttext.train_supervised("./train.txt", lr=1.0, epoch=50, wordNgrams=3)
result = classifier.test('./test.txt', k=-1,threshold=0.15)
print(result)
#
# jieba.load_userdict('./user_dict.txt')
#
# stopwords_set = set()
# with open('./stop_word.txt', 'r', encoding='utf-8') as infile:
#     for line in infile:
#         stopwords_set.add(line.strip())
#
# url = 'https://sunny-graduation.oss-cn-shanghai.aliyuncs.com/ppt/'
# with open('./ppt_map.txt', 'w', encoding='utf-8') as fout:
#     for root, dirs, files in os.walk(r".\ppt"):
#         for file in files:
#             prs = Presentation(os.path.join(root, file))
#             content = []
#             for _, slide in enumerate(prs.slides):
#                 for shape in slide.shapes:
#                     if shape.has_text_frame:
#                         text_frame = shape.text_frame
#                         for paragraph in text_frame.paragraphs:
#                             words = list(jieba.cut(paragraph.text))
#                             content.extend([item.strip() for item in words if
#                                             item.strip() not in stopwords_set and item.strip() != '' and not is_all_eng(
#                                                 item.strip())])
#
#             fout.write(file + ' ' + url + file + ' ' + ' '.join(classifier.predict(' '.join(content), k=5)[0]) + '\n')
